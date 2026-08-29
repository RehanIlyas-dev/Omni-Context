import uuid
from pathlib import Path
import chromadb
from docx import Document
from fastembed import TextEmbedding
from langchain_text_splitters import RecursiveCharacterTextSplitter
from pypdf import PdfReader  
from pptx import Presentation

BASE_DIR = Path(__file__).resolve().parent.parent


class DocumentIngestor:
    def __init__(self, db_path: str = None, collection_name: str = "omni_context"):
        self.db_path = db_path or str(BASE_DIR / "chroma_db")
        self.db_path = db_path
        self.collection_name = collection_name

        # Load the embedding model
        print("Loading embedding model (BAAI/bge-small-en-v1.5)...")
        self.embedding_model = TextEmbedding(model_name="BAAI/bge-small-en-v1.5")

        # Initialize ChromaDB client and collection
        print(f"Initializing ChromaDB at '{self.db_path}' and collection '{self.collection_name}'...")
        self.client = chromadb.PersistentClient(path=self.db_path)
        self.collection = self.client.get_or_create_collection(name=self.collection_name)

    def _load_txt_or_md(self, path: Path) -> str:
        with open(path, "r", encoding="utf-8") as f:
            return f.read()

    def _load_pdf(self, path: Path) -> str:
        with open(path, "rb") as f:
            head = f.read(5)
        if not head.startswith(b"%PDF"):
            raise ValueError(
                f"'{path.name}' is not a valid PDF (header: {head!r}). "
                "It may be an HTML/download page saved with a .pdf extension — re-download the actual PDF."
            )

        extracted_pages = []
        with PdfReader(str(path)) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    extracted_pages.append(text)
        return "\n\n".join(extracted_pages)

    def _load_pptx(self, path: Path) -> str:
        prs = Presentation(str(path))
        extracted_text = []
        for slide_idx, slide in enumerate(prs.slides):
            slide_content = [f"--- Slide {slide_idx + 1} ---"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            slide_content.append(paragraph.text.strip())
            extracted_text.append("\n".join(slide_content))
        return "\n\n".join(extracted_text)

    def _load_docx(self, path: Path) -> str:
       
        # Extract text from a DOCX file
        doc = Document(str(path))
        extracted_text = []

        # Extract body paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                extracted_text.append(para.text.strip())

        # Extract embedded table data
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_text:
                    extracted_text.append(" | ".join(row_text))

        return "\n\n".join(extracted_text)

    def load_document(self, file_path: str) -> str:
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        ext = path.suffix.lower()
        if ext in [".txt", ".md"]:
            return self._load_txt_or_md(path)
        elif ext == ".pdf":
            return self._load_pdf(path)
        elif ext in [".pptx", ".ppt"]:
            return self._load_pptx(path)
        elif ext in [".docx", ".doc"]:
            return self._load_docx(path)
        else:
            raise ValueError(f"Unsupported file format: {ext}")

    def chunk_text(self, text: str, chunk_size: int = 500, chunk_overlap: int = 100) -> list[str]:
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            separators=["\n\n", "\n", " ", ""]
        )
        return splitter.split_text(text)

    def ingest_file(self, file_path: str):
        path = Path(file_path)
        print(f"Ingesting file: {path.name}")

        content = self.load_document(file_path)
        if not content.strip():
            print(f"Warning: No text extracted from {path.name}")
            return

        chunks = self.chunk_text(content)
        print(f"Generated {len(chunks)} chunks.")

        embeddings = list(self.embedding_model.embed(chunks))

        ids = [str(uuid.uuid4()) for _ in chunks]
        metadatas = [
            {"source": path.name, "file_type": path.suffix.lower(), "chunk_index": i}
            for i in range(len(chunks))
        ]

        self.collection.add(
            ids=ids,
            documents=chunks,
            embeddings=embeddings,
            metadatas=metadatas
        )
        print(f"Successfully stored {len(chunks)} chunks from '{path.name}' in ChromaDB!\n")


if __name__ == "__main__":
    ingestor = DocumentIngestor()
    exts = ["*.txt", "*.md", "*.pdf", "*.pptx", "*.ppt", "*.docx", "*.doc"]
    files = sorted({p for ext in exts for p in Path("data").glob(ext)})
    if not files:
        print("No supported files found in data/")
    for file in files:
        try:
            ingestor.ingest_file(str(file))
        except Exception as e:
            print(f"Skipping '{file.name}': {e}\n")